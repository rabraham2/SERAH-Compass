from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import date

BASE = Path("Dataset/model_outputs")
png_dir = BASE / "onepagers"
pngs = sorted(png_dir.glob("onepager_*.png"))
out_ppt = BASE / "category_onepagers.pptx"

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)  # 16:9

# Cover
slide = prs.slides.add_slide(prs.slide_layouts[6])
tx = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.7), Inches(1.6)).text_frame
tx.text = "Category One-Pagers"
tx.paragraphs[0].font.size = Pt(48)
sub = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.7), Inches(1.6)).text_frame
sub.text = f"Baseline vs proposed price — {date.today():%d %b %Y}"

# Pages
for img in pngs:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(img), Inches(0.25), Inches(0.25),
                             width=Inches(12.83), height=Inches(7.0))
print("Saved:", out_ppt)
